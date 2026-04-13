#!/usr/bin/env python3
"""
Create or update a paper-reading PPT from a structured outline JSON.

The generated layout is fixed:
- white background
- top one-third text area
- bottom two-thirds image area
"""

from __future__ import annotations

import argparse
import re
import json
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Sequence, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt

try:
    from PIL import Image
except Exception:  # pragma: no cover - optional dependency in runtime
    Image = None

# ---------------------------------------------------------------------------
# 4-level section structure
# ---------------------------------------------------------------------------
SECTION_ORDER = ["motivation", "innovation", "methods", "results"]
SECTION_LABELS = {
    "motivation": "Motivation",
    "innovation": "Innovation",
    "methods":    "Methods",
    "results":    "Results",
}
SECTION_COLORS = {
    "motivation": RGBColor(180, 30,  30),   # dark red
    "innovation": RGBColor(30,  80,  180),  # dark blue
    "methods":    RGBColor(20,  130, 70),   # dark green
    "results":    RGBColor(130, 30,  150),  # purple
}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build paper reading pptx from JSON outline.")
    parser.add_argument("--outline-json", required=True, help="Path to outline JSON.")
    parser.add_argument(
        "--output",
        required=True,
        help="Output PPTX path or output directory. Final filename uses '<YYMMDD>_<title words>.pptx'.",
    )
    parser.add_argument("--base-pptx", default="", help="Existing PPTX to append/update.")
    parser.add_argument(
        "--date-tag",
        default="",
        help="Date tag used in final filename. Default: current date in YYMMDD.",
    )
    parser.add_argument(
        "--title-words",
        type=int,
        default=8,
        help="How many leading title words to keep in final filename (default: 8).",
    )
    parser.add_argument(
        "--keep-output-name",
        action="store_true",
        help="Keep the exact --output filename without automatic renaming.",
    )
    parser.add_argument(
        "--clear-existing",
        action="store_true",
        help="When base PPTX is provided, remove existing slides before adding new slides.",
    )
    parser.add_argument(
        "--force-16x9",
        action="store_true",
        help="Force 16:9 size even when using base PPTX.",
    )
    return parser.parse_args()


def remove_all_slides(prs: Presentation) -> None:
    # python-pptx has no public remove API. This method is standard practice.
    slide_ids = list(prs.slides._sldIdLst)  # pylint: disable=protected-access
    for slide_id in slide_ids:
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)  # pylint: disable=protected-access


def load_outline(path: Path) -> Dict[str, Any]:
    with path.open("r", encoding="utf-8") as f:
        data = json.load(f)
    if "slides" not in data or not isinstance(data["slides"], list):
        raise ValueError("outline JSON must include a 'slides' list.")
    return data


def _slugify_title_words(title: str, max_words: int) -> str:
    tokens = re.findall(r"[A-Za-z0-9]+", title)
    words = [tok for tok in tokens if tok]
    if max_words > 0:
        words = words[:max_words]
    slug = "_".join(words)
    slug = re.sub(r"_+", "_", slug).strip("_")
    return slug or "paper"


def infer_paper_title(outline: Dict[str, Any]) -> str:
    paper = outline.get("paper", {})
    if isinstance(paper, dict):
        title = str(paper.get("title", "")).strip()
        if title:
            return title

    slides = outline.get("slides", [])
    if isinstance(slides, list) and slides:
        first = slides[0]
        if isinstance(first, dict):
            title = str(first.get("title", "")).strip()
            if title:
                return title
    return "paper"


def resolve_output_path(
    requested_output: Path,
    outline: Dict[str, Any],
    date_tag: str,
    title_words: int,
) -> Path:
    parent = requested_output if requested_output.suffix.lower() != ".pptx" else requested_output.parent
    title = infer_paper_title(outline)
    safe_words = max(1, title_words)
    stem = f"{date_tag}_{_slugify_title_words(title, safe_words)}"
    return (parent / f"{stem}.pptx").resolve()


def ensure_white_background(slide: Any) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)


def style_title(paragraph: Any) -> None:
    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    if not paragraph.runs:
        return
    run = paragraph.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 0, 0)


def style_body(paragraph: Any) -> None:
    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    for run in paragraph.runs:
        run.font.name = "Calibri"
        run.font.size = Pt(16)
        run.font.bold = False
        run.font.color.rgb = RGBColor(0, 0, 0)


def add_text_block(
    slide: Any,
    title: str,
    sections: Dict[str, str],
    text_lines: Sequence[str],
) -> None:
    """Render the text area at the top of the slide.

    If *sections* contains any of the four keys (motivation / innovation /
    methods / results), those are rendered as labelled rows with a coloured
    bold label followed by the content.  Otherwise *text_lines* is used as a
    plain-text fallback (backwards compatible with old outlines).
    """
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(12.4), Inches(2.35))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    # --- Title row ---
    title_p = tf.paragraphs[0]
    title_p.text = title
    style_title(title_p)

    # --- 4-level sections or plain text_lines fallback ---
    has_sections = any(sections.get(k) for k in SECTION_ORDER)
    if has_sections:
        for key in SECTION_ORDER:
            content = sections.get(key, "").strip()
            if not content:
                continue
            p = tf.add_paragraph()
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            # Coloured bold label
            label_run = p.add_run()
            label_run.text = f"{SECTION_LABELS[key]}:  "
            label_run.font.name = "Calibri"
            label_run.font.size = Pt(13)
            label_run.font.bold = True
            label_run.font.color.rgb = SECTION_COLORS[key]
            # Plain content
            content_run = p.add_run()
            content_run.text = content
            content_run.font.name = "Calibri"
            content_run.font.size = Pt(13)
            content_run.font.bold = False
            content_run.font.color.rgb = RGBColor(30, 30, 30)
    else:
        for line in text_lines:
            p = tf.add_paragraph()
            p.text = line
            style_body(p)


def image_boxes(count: int) -> List[Tuple[float, float, float, float]]:
    count = max(1, min(count, 3))
    left = 0.45
    top = 2.75     # pushed down to accommodate taller text box
    area_w = 12.4
    area_h = 4.10  # total image height (bottom stays at ~6.85")
    gap = 0.20

    if count == 1:
        return [(left, top, area_w, area_h)]
    if count == 2:
        w = (area_w - gap) / 2
        return [(left, top, w, area_h), (left + w + gap, top, w, area_h)]
    w = (area_w - 2 * gap) / 3
    return [
        (left, top, w, area_h),
        (left + w + gap, top, w, area_h),
        (left + 2 * (w + gap), top, w, area_h),
    ]


def resolve_path(raw_path: str, outline_dir: Path) -> Path:
    p = Path(raw_path)
    return p if p.is_absolute() else (outline_dir / p).resolve()


def fit_size(path: Path, box_w: float, box_h: float) -> Tuple[float, float]:
    if Image is None:
        return box_w, box_h
    try:
        with Image.open(path) as img:
            iw, ih = img.size
        if ih <= 0 or iw <= 0:
            return box_w, box_h
        image_ratio = iw / ih
        box_ratio = box_w / box_h
        if image_ratio > box_ratio:
            w = box_w
            h = box_w / image_ratio
        else:
            h = box_h
            w = box_h * image_ratio
        return w, h
    except Exception:
        return box_w, box_h


def add_caption(slide: Any, text: str, x: float, y: float, w: float) -> None:
    cap = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.20))
    tf = cap.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    if p.runs:
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(60, 60, 60)


def add_image_block(
    slide: Any,
    image_paths: Sequence[Path],
    figure_refs: Sequence[str],
) -> None:
    existing = [p for p in image_paths if p.exists()][:3]
    if not existing:
        placeholder = slide.shapes.add_textbox(Inches(0.55), Inches(3.20), Inches(12.0), Inches(0.8))
        tf = placeholder.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "未提供可用配图（请在 outline.json 的 images 中填写 PDF 截图路径）。"
        if p.runs:
            p.runs[0].font.name = "Calibri"
            p.runs[0].font.size = Pt(15)
            p.runs[0].font.color.rgb = RGBColor(120, 120, 120)
        return

    boxes = image_boxes(len(existing))
    for idx, (img_path, box) in enumerate(zip(existing, boxes)):
        bx, by, bw, bh = box
        fw, fh = fit_size(img_path, bw, bh)
        x = bx + (bw - fw) / 2
        y = by + (bh - fh) / 2
        slide.shapes.add_picture(str(img_path), Inches(x), Inches(y), Inches(fw), Inches(fh))

        if idx < len(figure_refs):
            cap = figure_refs[idx]
        else:
            cap = img_path.stem
        add_caption(slide, cap, bx, by + bh + 0.03, bw)


def normalize_slide(slide_spec: Dict[str, Any]) -> Dict[str, Any]:
    title = str(slide_spec.get("title", "")).strip() or "未命名页面"

    # --- 4-level sections (new primary format) ---
    sections: Dict[str, str] = {}
    for key in SECTION_ORDER:
        val = str(slide_spec.get(key, "")).strip()
        if val:
            sections[key] = val

    # --- text_lines (legacy / fallback) ---
    text_lines = slide_spec.get("text_lines", [])
    if not isinstance(text_lines, list):
        text_lines = [str(text_lines)]
    text_lines = [str(line).strip() for line in text_lines if str(line).strip()]

    images = slide_spec.get("images", [])
    if not isinstance(images, list):
        images = [str(images)]
    images = [str(item).strip() for item in images if str(item).strip()]

    figure_refs = slide_spec.get("figure_refs", [])
    if not isinstance(figure_refs, list):
        figure_refs = [str(figure_refs)]
    figure_refs = [str(item).strip() for item in figure_refs if str(item).strip()]

    notes = str(slide_spec.get("notes", "")).strip()
    return {
        "title": title,
        "sections": sections,
        "text_lines": text_lines,
        "images": images,
        "figure_refs": figure_refs,
        "notes": notes,
    }


def build_presentation(
    outline: Dict[str, Any],
    outline_dir: Path,
    base_pptx: Path | None,
    clear_existing: bool,
    force_16x9: bool,
) -> Presentation:
    if base_pptx and base_pptx.exists():
        prs = Presentation(str(base_pptx))
        if clear_existing:
            remove_all_slides(prs)
    else:
        prs = Presentation()

    if force_16x9 or not (base_pptx and base_pptx.exists()):
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slides = [normalize_slide(spec) for spec in outline["slides"]]

    for spec in slides:
        slide = prs.slides.add_slide(layout)
        ensure_white_background(slide)
        add_text_block(slide, spec["title"], spec["sections"], spec["text_lines"])
        resolved_images = [resolve_path(raw, outline_dir) for raw in spec["images"]]
        add_image_block(slide, resolved_images, spec["figure_refs"])

        if spec["notes"]:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = spec["notes"]

    return prs


def main() -> None:
    args = parse_args()
    outline_path = Path(args.outline_json).resolve()
    requested_output = Path(args.output).resolve()
    base_path = Path(args.base_pptx).resolve() if args.base_pptx else None

    if not outline_path.exists():
        raise FileNotFoundError(f"Outline JSON not found: {outline_path}")

    outline = load_outline(outline_path)
    date_tag = args.date_tag.strip() or datetime.now().strftime("%y%m%d")
    if args.keep_output_name:
        output_path = requested_output
    else:
        output_path = resolve_output_path(
            requested_output=requested_output,
            outline=outline,
            date_tag=date_tag,
            title_words=args.title_words,
        )

    prs = build_presentation(
        outline=outline,
        outline_dir=outline_path.parent,
        base_pptx=base_path,
        clear_existing=args.clear_existing,
        force_16x9=args.force_16x9,
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"[OK] PPT generated: {output_path}")
    if not args.keep_output_name:
        print(f"[INFO] output naming rule: {date_tag}_<paper_title_words>.pptx")
    print(f"[INFO] slides count: {len(prs.slides)}")


if __name__ == "__main__":
    main()
